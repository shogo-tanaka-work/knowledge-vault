"""政府系シンクタンク風スライド: AI提案書作成ワークフローの現状"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

def _hex(h):
    return RGBColor((h >> 16) & 0xFF, (h >> 8) & 0xFF, h & 0xFF)

GREEN_MAIN = _hex(0x6FA84B)
GREEN_DARK = _hex(0x5B9B3E)
GREEN_LIGHT = _hex(0xE8F1DA)
GREEN_BORDER = _hex(0x8FBF5E)
GRAY_LABEL = _hex(0xBFBFBF)
GRAY_TEXT = _hex(0x595959)
TEXT_DARK = _hex(0x333333)
WHITE = _hex(0xFFFFFF)
BORDER_GRAY = _hex(0xA6A6A6)

FONT_JP = "Yu Gothic"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, *, size=10, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None,
             line_color=None, line_width=0.75, font=FONT_JP):
    tb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    tb.shadow.inherit = False
    if fill is None:
        tb.fill.background()
    else:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if line is False:
        tb.line.fill.background()
    elif line_color is not None:
        tb.line.color.rgb = line_color
        tb.line.width = Pt(line_width)
    else:
        tb.line.fill.background()
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def set_dash(shape, dash="dash"):
    ln = shape.line._get_or_add_ln()
    for tag in ("a:prstDash",):
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", dash)


# ============== タイトル ==============
add_text(slide, Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.55),
         "1-3 AI提案書作成ワークフローの現状", size=24, bold=True, color=TEXT_DARK)

# 緑の四角マーカー
marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.92),
                                Inches(0.12), Inches(0.12))
marker.fill.solid(); marker.fill.fore_color.rgb = GREEN_DARK
marker.line.fill.background(); marker.shadow.inherit = False

lead = (
    "本ワークフローは、以下の4つのStepで構成される。Step1「リード獲得」、Step2「ヒアリング」、Step3「提案書作成」、Step4「クロージング」の順に進行する。\n"
    "AI活用領域と人手作業領域が混在しており、特にStep1のリードスコアリングおよびStep3の提案書ドラフト生成ではAIによる効率化余地が大きい。一方、Step2のヒアリングおよび\n"
    "Step4のクロージングは人的関係構築・交渉が中核を担い、人手依存度が高い状態にある。本資料では各ステップにおけるAI活用度合いと残存課題を整理し、今後の効率化打ち手の\n"
    "検討材料とする。"
)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12.4), Inches(0.95),
         lead, size=10, color=TEXT_DARK)

# ============== 右上注記 ==============
note_x, note_y = Inches(11.3), Inches(1.95)
add_text(slide, note_x, note_y, Inches(1.85), Inches(0.22),
         "本資料のセクション番号と対応", size=9, color=GRAY_TEXT, align=PP_ALIGN.LEFT)

# ============== 4ステップ シェブロン ==============
arrow_y = Inches(2.22)
arrow_h = Inches(0.5)
arrow_w = Inches(2.95)
gap = Inches(0.05)
start_x = Inches(0.7)

steps = [
    ("Step1", "リード獲得"),
    ("Step2", "ヒアリング"),
    ("Step3", "提案書作成"),
    ("Step4", "クロージング"),
]
step_centers = []
for i, (step, label) in enumerate(steps):
    x = start_x + (arrow_w + gap) * i
    sh = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, x, arrow_y, arrow_w, arrow_h)
    sh.fill.solid(); sh.fill.fore_color.rgb = GREEN_MAIN
    sh.line.color.rgb = GREEN_MAIN
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.3)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = step
    r.font.name = FONT_JP; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.name = FONT_JP; r2.font.size = Pt(12); r2.font.bold = True; r2.font.color.rgb = WHITE
    step_centers.append(x + arrow_w / 2)

# ステップ間の縦点線
dash_top = Inches(2.85)
dash_bottom = Inches(7.0)
for i in range(1, 4):
    line_x = start_x + (arrow_w + gap) * i - gap / 2
    ln = slide.shapes.add_connector(1, line_x, dash_top, line_x, dash_bottom)
    ln.line.color.rgb = GRAY_LABEL
    ln.line.width = Pt(0.75)
    set_dash(ln, "sysDash")

# ============== 行ラベル ==============
row_label_x = Inches(0.4)
row_label_w = Inches(0.55)
ai_y = Inches(3.35); ai_h = Inches(1.7)
human_y = Inches(5.35); human_h = Inches(1.6)

for (y, h, txt) in [(ai_y, ai_h, "AI活用\n領域"), (human_y, human_h, "人手作業\n領域")]:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, row_label_x, y, row_label_w, h)
    box.fill.solid(); box.fill.fore_color.rgb = GRAY_LABEL
    box.line.fill.background(); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    for i, ln in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.name = FONT_JP; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE


def chapter_box(x, y, w, header_text, body_title, sub_boxes, *, body_h_in=1.45):
    """章ヘッダ＋破線エリア＋サブボックス"""
    hdr_h = Inches(0.55)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, hdr_h)
    hdr.fill.solid(); hdr.fill.fore_color.rgb = GREEN_DARK
    hdr.line.fill.background(); hdr.shadow.inherit = False
    tf = hdr.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = header_text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.name = FONT_JP; r.font.bold = True; r.font.color.rgb = WHITE
        r.font.size = Pt(9) if i == 0 else Pt(10)

    # 破線エリア
    body_y = y + hdr_h
    body_h = Inches(body_h_in)
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, body_y, w, body_h)
    body.fill.solid(); body.fill.fore_color.rgb = GREEN_LIGHT
    body.line.color.rgb = GREEN_BORDER; body.line.width = Pt(1.25)
    set_dash(body, "dash")
    body.shadow.inherit = False
    tf = body.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = body_title
    r.font.name = FONT_JP; r.font.size = Pt(9); r.font.bold = False; r.font.color.rgb = TEXT_DARK

    # サブボックス
    n = len(sub_boxes)
    if n == 0:
        return
    inner_pad = Inches(0.12)
    sub_top = body_y + Inches(0.42)
    sub_h = Inches(0.55)
    if n == 1:
        sw = w - inner_pad * 2
        sb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + inner_pad, sub_top, sw, sub_h)
        _style_sub(sb, sub_boxes[0])
    else:
        # 2個: 左右に並べ、間に矢印
        arrow_w_in = Inches(0.25)
        sw = (w - inner_pad * 2 - arrow_w_in) / 2
        sb1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + inner_pad, sub_top, sw, sub_h)
        _style_sub(sb1, sub_boxes[0])
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + inner_pad + sw, sub_top + Inches(0.15),
                                       arrow_w_in, Inches(0.25))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = GRAY_TEXT
        arrow.line.fill.background(); arrow.shadow.inherit = False
        sb2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x + inner_pad + sw + arrow_w_in, sub_top, sw, sub_h)
        _style_sub(sb2, sub_boxes[1])


def _style_sub(shape, text):
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER_GRAY; shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.name = FONT_JP; r.font.size = Pt(9); r.font.color.rgb = TEXT_DARK


# 各ステップのコンテンツエリアX座標 (シェブロンの直下に揃える)
step_x = [start_x + (arrow_w + gap) * i for i in range(4)]
content_pad = Inches(0.15)
content_w = arrow_w - content_pad * 2

# ============== Step1: 2章 (上段単独) ==============
chapter_box(step_x[0] + content_pad, Inches(2.95), content_w,
            "2章\n現状リード獲得状況の把握",
            "本ワークフローの起点となるリード獲得チャネルにつき確認",
            ["自社の主要リードチャネル・案件発生件数の確認"],
            body_h_in=1.3)

# ============== Step4: 8章 (上段単独) ==============
chapter_box(step_x[3] + content_pad, Inches(2.95), content_w,
            "8章\n提案・クロージング結果",
            "提案実施および受注／失注結果の振り返り",
            ["提案実施・受注／失注の振り返り"],
            body_h_in=1.3)

# ============== AI活用領域 (Step2 × 2章, Step3 × 1章) ==============
# Step2 上: 3章
chapter_box(step_x[1] + content_pad, Inches(3.30), content_w,
            "3章\n案件データの分析",
            "CRM上の蓄積データを用いた優先案件の特定",
            ["CRM上の案件\n履歴整理", "リードスコアリング・優先度判定"],
            body_h_in=1.0)
# Step2 下: 4章
chapter_box(step_x[1] + content_pad, Inches(4.85+0.1), content_w,
            "4章\nヒアリング前リサーチ",
            "AIによる業界・競合情報の自動収集と仮説生成",
            ["業界・競合情報の\n自動収集", "想定課題仮説\nの生成"],
            body_h_in=1.0)

# Step3 AI: 7章
chapter_box(step_x[2] + content_pad, Inches(3.30), content_w,
            "7章\n過去提案資産との比較",
            "類似案件のRAG検索および提案書ドラフトの自動生成",
            ["類似案件のRAG\n検索", "提案書ドラフト\nの自動生成"],
            body_h_in=1.0)

# ============== 人手作業領域 (Step2 × 5章, Step3 × 6章) ==============
# 5章: Step2 人手
chapter_box(step_x[1] + content_pad, Inches(5.55), content_w,
            "5章\n顧客ヒアリング実施",
            "顧客課題の深掘りと関係構築",
            ["対面/オンラインによる課題深掘りと関係構築"],
            body_h_in=1.1)

# 6章: Step3 人手
chapter_box(step_x[2] + content_pad, Inches(5.55), content_w,
            "6章\n提案書のカスタマイズ・レビュー",
            "顧客文脈反映および社内レビューの実施",
            ["顧客文脈の反映と社内レビュー"],
            body_h_in=1.1)

# ============== フッター ==============
add_text(slide, Inches(0.4), Inches(7.18), Inches(8), Inches(0.25),
         "AI提案書作成ワークフロー 現状分析 2026", size=9, color=GRAY_TEXT)
add_text(slide, Inches(8.5), Inches(7.18), Inches(4.6), Inches(0.25),
         "© 2026. Internal use only.", size=9, color=GRAY_TEXT, align=PP_ALIGN.RIGHT)
add_text(slide, Inches(0.4), Inches(7.18), Inches(0.3), Inches(0.25),
         "6", size=9, color=GRAY_TEXT)

out = Path(__file__).resolve().parent.parent / "output" / "ai_proposal_workflow_current.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"saved: {out}")
