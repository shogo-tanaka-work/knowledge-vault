#!/usr/bin/env python3
"""
tmp-source/12-ai-workflow-asis.html の内容を、編集可能な PowerPoint (.pptx) として再構築する。
- 文字はすべてネイティブテキストボックス
- 装飾はすべて PowerPoint シェイプ
- 画像埋め込み禁止（DB/書類アイコンは MSO_SHAPE で近似）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ---------- 単位変換 ----------
PX_EMU = 9525  # 1280 px = 12,192,000 EMU = 13.333 inch
def E(px): return int(px * PX_EMU)
def P(px): return Pt(px * 0.75)

# ---------- カラーパレット (slide 12 と同じ) ----------
NAVY        = RGBColor(0x1A, 0x2B, 0x4A)
BRONZE      = RGBColor(0xB8, 0x94, 0x56)
GRAY_BAND   = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_BORDER = RGBColor(0xC9, 0xCD, 0xD1)
HIGHLIGHT   = RGBColor(0xD6, 0xE4, 0xF5)
TEXT_BODY   = RGBColor(0x33, 0x33, 0x33)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_777    = RGBColor(0x77, 0x77, 0x77)
GRAY_444    = RGBColor(0x44, 0x44, 0x44)

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
def qn(tag): return '{%s}%s' % (NS_A, tag)

# ---------- 基本ヘルパ ----------
def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=1.2, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, E(x), E(y), E(w), E(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Emu(int(line_w * PX_EMU))
    # 角丸の半径を小さく
    if rounded:
        s.adjustments[0] = 0.06
    s.shadow.inherit = False
    return s

def fmt_text(shape, text, size_px=11, bold=False, color=TEXT_BODY,
             align='center', valign='middle', font='Noto Sans JP'):
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Emu(int(3 * PX_EMU))
    tf.margin_top = tf.margin_bottom = Emu(int(1 * PX_EMU))
    tf.word_wrap = True
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}[valign]
    # 1段落目をテキストにし、改行があれば追加段落
    lines = text.split('\n')
    tf.text = lines[0]
    for ln in lines[1:]:
        p = tf.add_paragraph()
        p.text = ln
    for p in tf.paragraphs:
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
        for run in p.runs:
            run.font.size = P(size_px)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font

def add_textbox(slide, x, y, w, h, text, **kw):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    fmt_text(tb, text, **kw)
    return tb

def add_arrow(slide, x1, y1, x2, y2, color=GRAY_444, line_w=1.0):
    """矢印つき直線（始点→終点に三角ヘッド）。
       python-pptx の add_connector は (begin_x, begin_y, end_x, end_y) の絶対座標を取る。"""
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2))
    cn.line.color.rgb = color
    cn.line.width = Emu(int(line_w * PX_EMU))
    # 矢印ヘッドを XML で追加
    ln = cn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'sm')
    tail.set('len', 'sm')
    return cn

def add_box_with_text(slide, x, y, w, h, text, fill=WHITE, line=NAVY, line_w=1.2,
                     size_px=10.5, color=NAVY, bold=False, rounded=True, align='center'):
    s = add_rect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w, rounded=rounded)
    fmt_text(s, text, size_px=size_px, color=color, bold=bold, align=align)
    return s

def add_num_badge(slide, cx, cy, n, size=18):
    """中心 (cx, cy) に直径 size の濃紺丸＋白抜き番号"""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, E(cx - size/2), E(cy - size/2), E(size), E(size))
    s.fill.solid(); s.fill.fore_color.rgb = NAVY
    s.line.fill.background()
    s.shadow.inherit = False
    fmt_text(s, str(n), size_px=12, color=WHITE, bold=True, align='center', valign='middle')
    return s

def add_db(slide, x, y, w, h, label):
    """DB シリンダー（drum）＋ラベルを内側に"""
    s = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_MAGNETIC_DISK, E(x), E(y), E(w), E(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = NAVY; s.line.width = Emu(int(1.2 * PX_EMU))
    s.shadow.inherit = False
    fmt_text(s, label, size_px=10, color=NAVY, bold=True, align='center', valign='middle')
    return s

def add_doc(slide, x, y, w, h):
    """書類アイコン（折り返し角つき）"""
    s = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER, E(x), E(y), E(w), E(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = NAVY; s.line.width = Emu(int(1.2 * PX_EMU))
    s.shadow.inherit = False
    return s

def add_chevron(slide, x, y, w, h, text, fill=BRONZE, color=WHITE, size_px=11.5, bold=True, notch_left=False):
    tip = min(48, w * 0.12)
    vertices = []
    if notch_left:
        vertices.append((E(x), E(y)))
        vertices.append((E(x + tip), E(y + h / 2)))
        vertices.append((E(x), E(y + h)))
    else:
        vertices.append((E(x), E(y)))
        vertices.append((E(x), E(y + h)))
    vertices.extend([
        (E(x + w - tip), E(y + h)),
        (E(x + w), E(y + h / 2)),
        (E(x + w - tip), E(y)),
    ])
    builder = slide.shapes.build_freeform(vertices[0][0], vertices[0][1])
    builder.add_line_segments(vertices[1:], close=True)
    s = builder.convert_to_shape()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    fmt_text(s, text, size_px=size_px, color=color, bold=bold, align='center', valign='middle')
    return s

def add_tag(slide, x, y, w, h, text):
    """ブロンズ枠 + 白背景の小さなタグ（画面/タスク/マクロ）"""
    s = add_rect(slide, x, y, w, h, fill=WHITE, line=BRONZE, line_w=1.2, rounded=False)
    fmt_text(s, text, size_px=9.5, color=TEXT_BODY, align='center', valign='middle')
    return s

# =====================================================================
# メイン
# =====================================================================
def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ==========================================================
    # 1) ヘッダ部（章番号・サブタイトル・アクションタイトル）
    # ==========================================================
    add_textbox(slide, 44, 28, 200, 30, "2.1.",
                size_px=22, bold=True, color=NAVY, align='left', valign='top')
    add_textbox(slide, 44, 58, 800, 18, "2.1.3. AI 提案書作成ワークフローの現状 〜 提示まで",
                size_px=13, color=TEXT_BODY, align='left', valign='top')

    # アクションタイトル：左にブロンズ縦バー＋テキスト
    add_rect(slide, 44, 82, 5, 50, fill=BRONZE, line=None)
    add_textbox(slide, 56, 82, 850, 50,
                "提案書作成までの全 9 工程のうち 8 工程を AI ツール群が担い、\n"
                "営業担当は最終確認 1 工程のみ担当している。",
                size_px=19, bold=True, color=NAVY, align='left', valign='top')

    # ==========================================================
    # 2) 凡例ボックス（右上）
    # ==========================================================
    LEG_X, LEG_Y, LEG_W, LEG_H = 920, 28, 320, 76
    add_rect(slide, LEG_X, LEG_Y, LEG_W, LEG_H, fill=WHITE, line=GRAY_BORDER, line_w=1.0)
    add_textbox(slide, LEG_X + 8, LEG_Y + 4, 100, 16, "【凡例】",
                size_px=11, bold=True, color=TEXT_BODY, align='left', valign='top')
    # 3 chips
    chip_y = LEG_Y + 24
    chip_h = 40
    chip_pad = 8
    chip_w = (LEG_W - chip_pad * 4) / 3
    cx = LEG_X + chip_pad
    # gray
    s = add_rect(slide, cx, chip_y, chip_w, chip_h, fill=GRAY_BAND, line=None)
    fmt_text(s, "機能群", size_px=10.5, bold=True, color=TEXT_BODY)
    cx += chip_w + chip_pad
    # bronze outline
    s = add_rect(slide, cx, chip_y, chip_w, chip_h, fill=WHITE, line=BRONZE, line_w=1.5)
    fmt_text(s, "実装方式", size_px=10.5, bold=True, color=TEXT_BODY)
    cx += chip_w + chip_pad
    # navy fill
    s = add_rect(slide, cx, chip_y, chip_w, chip_h, fill=NAVY, line=None)
    fmt_text(s, "運用事業者\n独自開発の機能", size_px=10, bold=True, color=WHITE)

    # ==========================================================
    # 3) リード文（●×2）
    # ==========================================================
    add_textbox(slide, 50, 144, 1180, 24,
                "●  AI ツール群は登録した商談データに基づき顧客 DB・案件 DB の更新及び CSV 出力を行い、最終成果物である提案書一式の作成を実施する。",
                size_px=11.5, color=TEXT_BODY, align='left', valign='top')
    add_textbox(slide, 50, 164, 1180, 24,
                "●  営業担当は AI ツール群から受領した提案書一式を確認のうえクライアント商談で提示する。",
                size_px=11.5, color=TEXT_BODY, align='left', valign='top')

    # ==========================================================
    # 4) アクター + ブロンズシェブロンバンド
    # ==========================================================
    # アクター 2 つ（左：AI ツール群、右：営業担当）
    add_textbox(slide, 50, 196, 460, 18, "👤 AI ツール群（運用基盤）",
                size_px=11, bold=True, color=NAVY, align='center', valign='middle')
    add_textbox(slide, 510, 196, 460, 18, "👤 営業担当（事業部門）",
                size_px=11, bold=True, color=NAVY, align='center', valign='middle')

    # ブロンズ→ネイビーのシェブロンバンド
    band_y = 220
    band_h = 30
    add_chevron(slide, 50,  band_y, 460, band_h, "受領・集計（AI ツール群）",
                fill=BRONZE, color=WHITE)
    add_chevron(slide, 510, band_y, 460, band_h, "提示（営業担当）",
                fill=NAVY, color=WHITE, notch_left=True)

    # ==========================================================
    # 5) 図解枠（点線・グレー）
    # ==========================================================
    diag_x, diag_y, diag_w, diag_h = 50, 256, 920, 360
    diag_box = add_rect(slide, diag_x, diag_y, diag_w, diag_h, fill=WHITE, line=GRAY_BORDER, line_w=1.0)
    # 点線にする（XML で dashStyle 指定）
    ln = diag_box.line._get_or_add_ln()
    prstDash = etree.SubElement(ln, qn('prstDash'))
    prstDash.set('val', 'dash')

    # 図解内部の座標系：HTML SVG viewBox 0..940 を diag_w にマップ
    SX = diag_x + 8       # 左マージン
    SY = diag_y + 8       # 上マージン
    SCALE = (diag_w - 16) / 940.0    # ≈ 0.962
    def sx(px): return SX + px * SCALE
    def sy(py): return SY + py * SCALE
    def sw(pw): return pw * SCALE
    def sh(ph): return ph * SCALE

    # ==========================================================
    # 6) レーン1（旧 SFA システム）
    # ==========================================================
    # システム名（左ラベル）
    add_textbox(slide, sx(2), sy(20), sw(70), sh(40), "旧 SFA\nシステム",
                size_px=11, bold=True, color=NAVY, align='center', valign='middle')

    # ボックス: 取込・顧客DB更新
    add_box_with_text(slide, sx(98), sy(50), sw(92), sh(36), "取込・\n顧客 DB 更新",
                     size_px=10.5, color=NAVY, line_w=1.2)
    add_tag(slide, sx(116), sy(92), sw(56), sh(14), "画面")
    add_num_badge(slide, sx(98), sy(50), 4)
    add_arrow(slide, sx(190), sy(68), sx(226), sy(68))

    # 顧客 DB
    add_db(slide, sx(226), sy(44), sw(56), sh(48), "顧客 DB")
    add_arrow(slide, sx(282), sy(68), sx(316), sy(68))

    # 案件サマリ更新
    add_box_with_text(slide, sx(316), sy(50), sw(92), sh(36), "案件サマリ更新\n（集計）",
                     size_px=10.5, color=NAVY)
    add_tag(slide, sx(334), sy(92), sw(56), sh(14), "画面")
    add_num_badge(slide, sx(316), sy(50), 5)
    add_arrow(slide, sx(408), sy(68), sx(444), sy(68))

    # 案件 DB
    add_db(slide, sx(444), sy(44), sw(56), sh(48), "案件 DB")
    add_arrow(slide, sx(500), sy(68), sx(534), sy(68))

    # データ加工・CSV
    add_box_with_text(slide, sx(534), sy(50), sw(92), sh(36), "データ加工・\nCSV 出力",
                     size_px=10.5, color=NAVY)
    add_tag(slide, sx(552), sy(92), sw(56), sh(14), "画面")

    # 旧 → エラーチェック (下に伸びる矢印)
    add_arrow(slide, sx(254), sy(92), sx(254), sy(138))
    add_box_with_text(slide, sx(180), sy(138), sw(148), sh(36),
                     "エラーチェック・疑義・\n報告状況等照会",
                     size_px=10.5, color=NAVY)
    add_tag(slide, sx(226), sy(180), sw(56), sh(14), "画面")
    add_num_badge(slide, sx(180), sy(138), 4)

    # チェック用提案書 (右、書類アイコン)
    add_doc(slide, sx(640), sy(50), sw(34), sh(44))
    add_textbox(slide, sx(630), sy(98), sw(54), sh(28),
                "チェック用\n提案書 (13 件)",
                size_px=9.5, color=TEXT_BODY, align='center', valign='top')
    add_arrow(slide, sx(626), sy(68), sx(640), sy(68))
    add_num_badge(slide, sx(640), sy(50), 7)

    # ==========================================================
    # 7) レーン2（新クラウド AI システム）
    # ==========================================================
    add_textbox(slide, sx(2), sy(206), sw(70), sh(40),
                "新クラウド\nAI システム",
                size_px=11, bold=True, color=NAVY, align='center', valign='middle')

    add_box_with_text(slide, sx(98), sy(230), sw(92), sh(36),
                     "取込・\n顧客 DB 更新", size_px=10.5, color=NAVY)
    add_tag(slide, sx(116), sy(272), sw(56), sh(14), "タスク")
    add_num_badge(slide, sx(98), sy(230), 4)
    add_arrow(slide, sx(190), sy(248), sx(226), sy(248))

    add_db(slide, sx(226), sy(224), sw(56), sh(48), "顧客 DB")
    add_arrow(slide, sx(282), sy(248), sx(316), sy(248))

    add_box_with_text(slide, sx(316), sy(230), sw(92), sh(36),
                     "案件サマリ更新\n（集計）", size_px=10.5, color=NAVY)
    add_tag(slide, sx(334), sy(272), sw(56), sh(14), "タスク")
    add_num_badge(slide, sx(316), sy(230), 5)
    add_arrow(slide, sx(408), sy(248), sx(444), sy(248))

    add_db(slide, sx(444), sy(224), sw(56), sh(48), "案件 DB")
    add_arrow(slide, sx(500), sy(248), sx(534), sy(248))

    add_box_with_text(slide, sx(534), sy(230), sw(92), sh(36),
                     "データ加工・\nCSV 出力", size_px=10.5, color=NAVY)
    add_tag(slide, sx(552), sy(272), sw(56), sh(14), "タスク")

    # レーン2 → エラーチェック (下)
    add_arrow(slide, sx(254), sy(272), sx(254), sy(304))
    add_box_with_text(slide, sx(180), sy(304), sw(148), sh(34),
                     "エラーチェック・疑義・\n報告状況等照会", size_px=10.5, color=NAVY)
    add_tag(slide, sx(226), sy(342), sw(56), sh(14), "画面")
    add_num_badge(slide, sx(180), sy(304), 4)

    # → チェック用データ参照 (highlight)
    add_arrow(slide, sx(328), sy(320), sx(358), sy(320))
    add_box_with_text(slide, sx(358), sy(304), sw(110), sh(34),
                     "チェック用\nデータ参照",
                     fill=HIGHLIGHT, size_px=10.5, color=NAVY, bold=True)
    add_tag(slide, sx(385), sy(342), sw(56), sh(14), "画面")
    add_num_badge(slide, sx(358), sy(304), 6)

    # CSV ファイル
    add_doc(slide, sx(612), sy(304), sw(30), sh(38))
    add_textbox(slide, sx(605), sy(346), sw(44), sh(14), "CSV",
                size_px=9.5, color=TEXT_BODY, align='center', valign='top')
    add_arrow(slide, sx(626), sy(290), sx(626), sy(304))

    # ==========================================================
    # 8) 提案書作成（中央右、マクロタグ）
    # ==========================================================
    add_box_with_text(slide, sx(704), sy(160), sw(80), sh(40),
                     "提案書作成", size_px=11, color=NAVY, bold=True, line_w=1.4)
    add_tag(slide, sx(720), sy(206), sw(48), sh(14), "マクロ")
    add_num_badge(slide, sx(704), sy(160), 8)

    # 7 → 帳票作成 / CSV → 帳票作成
    add_arrow(slide, sx(678), sy(78),  sx(704), sy(170))
    add_arrow(slide, sx(642), sy(320), sx(704), sy(190))

    # ==========================================================
    # 9) 出力帳票 3 つ
    # ==========================================================
    # 上：経営層提出資料
    add_doc(slide, sx(836), sy(48), sw(30), sh(38))
    add_textbox(slide, sx(810), sy(90), sw(80), sh(40),
                "経営層\n提出資料\n(73 件)",
                size_px=9.5, color=TEXT_BODY, align='center', valign='top')
    add_arrow(slide, sx(784), sy(170), sx(836), sy(68))
    add_num_badge(slide, sx(836), sy(48), 9)
    # 中：クライアント提出帳票
    add_doc(slide, sx(836), sy(160), sw(30), sh(38))
    add_textbox(slide, sx(810), sy(202), sw(80), sh(40),
                "クライアント\n提出帳票\n(51 件)",
                size_px=9.5, color=TEXT_BODY, align='center', valign='top')
    add_arrow(slide, sx(784), sy(180), sx(836), sy(180))
    add_num_badge(slide, sx(836), sy(160), 9)
    # 下：社内共有提案書
    add_doc(slide, sx(836), sy(276), sw(30), sh(38))
    add_textbox(slide, sx(810), sy(318), sw(80), sh(40),
                "社内共有\n提案書\n(1 件)",
                size_px=9.5, color=TEXT_BODY, align='center', valign='top')
    add_arrow(slide, sx(784), sy(195), sx(836), sy(295))
    add_num_badge(slide, sx(836), sy(276), 9)

    # ==========================================================
    # 10) 右パネル：主な業務内容
    # ==========================================================
    PX0, PY0 = 990, 130
    add_textbox(slide, PX0, PY0, 240, 22, "主な業務内容",
                size_px=14, bold=True, color=TEXT_BODY, align='left', valign='top')
    # 区切り線
    add_rect(slide, PX0, PY0 + 22, 240, 1.5, fill=GRAY_BORDER, line=None)

    work_items = [
        (4, "エラーチェック・疑義の確認及び修正",
            "提出された商談データにエラー・疑義があった場合、営業担当に確認の上修正を実施する（※現行業務では旧 SFA 側の画面を利用）"),
        (5, "集計・出力",
            "登録された商談データから、顧客 DB・案件 DB の更新を行い、データの加工及び CSV の出力を実施する"),
        (6, "チェック用データ参照",
            "AI ツール群にて独自に構築したダッシュボードから各種 DB のチェックを実施する／本画面からチェック用提案書を作成することも可能"),
        (7, "チェック用帳票作成",
            "新クラウド AI 側にはチェック用に利用可能な帳票出力機能がないため、旧 SFA 側の帳票を利用してチェックを実施する"),
        (8, "提案書作成",
            "主に旧 SFA システムのマクロを利用して、最終成果物である提案書を作成する（※現行で出力しているのは 1 帳票のみ）"),
        (9, "提示",
            "提案書一式を AI ツール群から営業担当に引き渡し、営業担当がクライアント面談で提示する"),
    ]
    iy = PY0 + 32
    item_h = 76
    for n, ttl, desc in work_items:
        add_num_badge(slide, PX0 + 10, iy + 6, n, size=20)
        add_textbox(slide, PX0 + 28, iy, 220, 18, ttl,
                    size_px=11.5, bold=True, color=TEXT_BODY, align='left', valign='top')
        add_textbox(slide, PX0 + 28, iy + 18, 220, item_h - 18, desc,
                    size_px=10, color=GRAY_444, align='left', valign='top')
        iy += item_h

    # ==========================================================
    # 11) KEY TAKEAWAY
    # ==========================================================
    KT_X, KT_Y, KT_W, KT_H = 50, 622, 920, 36
    add_rect(slide, KT_X, KT_Y, KT_W, KT_H, fill=GRAY_BAND, line=None)
    add_rect(slide, KT_X, KT_Y, 5, KT_H, fill=BRONZE, line=None)
    add_textbox(slide, KT_X + 14, KT_Y + 6, 130, 24, "KEY TAKEAWAY:",
                size_px=11.5, bold=True, color=NAVY, align='left', valign='middle')
    add_textbox(slide, KT_X + 144, KT_Y + 6, KT_W - 154, 24,
                "要するに、旧 SFA システムと新クラウド AI システムが並行稼働し、最終提案書作成は旧 SFA のマクロに依存している現状である。",
                size_px=11.5, color=TEXT_BODY, align='left', valign='middle')

    # ==========================================================
    # 12) フッタ
    # ==========================================================
    add_textbox(slide, 44, 678, 800, 26,
                "2.1. 業務に係る現状整理 ／ 2.1.3. AI 提案書作成ワークフローの現状 〜 提示まで\n出典：画像提供元情報をもとに作成",
                size_px=9.5, color=GRAY_777, align='left', valign='top')
    add_textbox(slide, 1180, 686, 60, 18, "12",
                size_px=11, bold=True, color=NAVY, align='right', valign='top')

    # ==========================================================
    # 保存
    # ==========================================================
    out = Path(__file__).resolve().parent.parent / "outputs" / "12-ai-workflow-asis-editable.pptx"
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"saved: {out}  ({out.stat().st_size/1024:.1f} KB)")
    return out


if __name__ == "__main__":
    build()
