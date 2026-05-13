#!/usr/bin/env python3
"""
政府系シンクタンク（Deloitte Tohmatsu）風スライドのレイアウトを再現し、
テーマを「AI 提案書作成ワークフローの現状」に置換した編集可能 .pptx を生成する。

- 4 Step バー（Step1 商談 / Step2 分析 / Step3 戦略 / Step4 提案）
- 2 行アプローチ（AI ツール担当 / 人手担当）
- 章番号バッジ付き白枠ボックス + 矢印
- 右側「本報告書の章番号と対応」縦パネル
- 文字はネイティブテキスト、装飾はすべて PowerPoint シェイプ
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ---------- 単位変換 (1280px 基準) ----------
PX_EMU = 9525
def E(px): return int(px * PX_EMU)
def P(px): return Pt(px * 0.75)

# ---------- カラーパレット ----------
NAVY        = RGBColor(0x1A, 0x2B, 0x4A)
BRONZE      = RGBColor(0xB8, 0x94, 0x56)
GRAY_BAR    = RGBColor(0xE8, 0xE8, 0xE8)   # Step バー
GRAY_ROW    = RGBColor(0xD9, 0xD9, 0xD9)   # 行ラベル帯
GRAY_BORDER = RGBColor(0xC9, 0xCD, 0xD1)
GRAY_777    = RGBColor(0x77, 0x77, 0x77)
GRAY_444    = RGBColor(0x44, 0x44, 0x44)
TEXT_BODY   = RGBColor(0x33, 0x33, 0x33)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
def qn(tag): return '{%s}%s' % (NS_A, tag)

# ---------- 基本ヘルパ ----------
def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=1.2, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, E(x), E(y), E(w), E(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Emu(int(line_w * PX_EMU))
    if rounded:
        s.adjustments[0] = 0.08
    s.shadow.inherit = False
    return s

def fmt_text(shape, text, size_px=11, bold=False, color=TEXT_BODY,
             align='center', valign='middle', font='Noto Sans JP'):
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Emu(int(3 * PX_EMU))
    tf.margin_top = tf.margin_bottom = Emu(int(1 * PX_EMU))
    tf.word_wrap = True
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}[valign]
    lines = text.split('\n')
    tf.text = lines[0]
    for ln in lines[1:]:
        p = tf.add_paragraph()
        p.text = ln
    for p in tf.paragraphs:
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
        for run in p.runs:
            run.font.size = P(size_px)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
            # eastAsia フォントも明示
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn('ea'))
            if ea is None:
                ea = etree.SubElement(rPr, qn('ea'))
            ea.set('typeface', font)
            latin = rPr.find(qn('latin'))
            if latin is None:
                latin = etree.SubElement(rPr, qn('latin'))
            latin.set('typeface', font)

def add_textbox(slide, x, y, w, h, text, **kw):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    fmt_text(tb, text, **kw)
    return tb

def add_arrow(slide, x1, y1, x2, y2, color=GRAY_444, line_w=1.0, head='triangle'):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2))
    cn.line.color.rgb = color
    cn.line.width = Emu(int(line_w * PX_EMU))
    ln = cn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('tailEnd'))
    tail.set('type', head)
    tail.set('w', 'sm')
    tail.set('len', 'sm')
    return cn

def add_box_with_text(slide, x, y, w, h, text, fill=WHITE, line=GRAY_BORDER, line_w=1.2,
                     size_px=10.5, color=NAVY, bold=False, rounded=False, align='center'):
    s = add_rect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w, rounded=rounded)
    fmt_text(s, text, size_px=size_px, color=color, bold=bold, align=align)
    return s

def add_num_badge(slide, cx, cy, n, size=22, fill=NAVY, color=WHITE, size_px=12):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, E(cx - size/2), E(cy - size/2), E(size), E(size))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    fmt_text(s, str(n), size_px=size_px, color=color, bold=True, align='center', valign='middle')
    return s

def add_chevron(slide, x, y, w, h, text, fill=NAVY, color=WHITE, size_px=12, bold=True):
    tip = min(28, w * 0.06)
    vertices = [
        (E(x), E(y)),
        (E(x), E(y + h)),
        (E(x + w - tip), E(y + h)),
        (E(x + w), E(y + h / 2)),
        (E(x + w - tip), E(y)),
    ]
    builder = slide.shapes.build_freeform(vertices[0][0], vertices[0][1])
    builder.add_line_segments(vertices[1:], close=True)
    s = builder.convert_to_shape()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    fmt_text(s, text, size_px=size_px, color=color, bold=bold, align='center', valign='middle')
    return s


# =====================================================================
# メイン
# =====================================================================
def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ==========================================================
    # 1) ヘッダ
    # ==========================================================
    add_textbox(slide, 44, 26, 800, 32, "1-3  AI 提案書作成ワークフローの現状",
                size_px=20, bold=True, color=NAVY, align='left', valign='top')
    # 下線
    add_rect(slide, 44, 60, 1180, 1.5, fill=NAVY, line=None)

    # ■リード文
    lead = (
        "■ 本ワークフローは、以下の 4Step で構成される。Step1 では「商談データの確認」、Step2 では「商談データの分析」を行い、"
        "Step3 では「提案戦略の検討」を実施する。\n"
        "    具体的には、Step1 は SFA から取得した商談データの正規化と整合性確認に集中する。"
        "Step2 は AI による顧客属性・案件履歴分析と、営業担当によるヒアリング整理を並行で進める。"
        "Step3 は AI が候補生成・営業が戦略選定を担い、最後に Step4 で本提案書を作成・提示する。"
    )
    add_textbox(slide, 44, 72, 1180, 56, lead,
                size_px=10.5, color=TEXT_BODY, align='left', valign='top')

    # 右上「本報告書の章番号と対応」ラベル
    add_textbox(slide, 1080, 132, 160, 18, "本報告書の章番号と対応",
                size_px=10, bold=True, color=NAVY, align='center', valign='middle')

    # ==========================================================
    # 2) メイングリッド領域定義
    # ==========================================================
    # キャンバス: 1280 x 750
    # 左列: 行ラベル (90)
    # 中央 4 Step 列: それぞれ約 220
    # 右列: 章対応パネル (約 180)
    GRID_X = 50
    GRID_Y = 158
    LABEL_W = 90
    STEP_W = 215
    PANEL_W = 200
    PANEL_X = GRID_X + LABEL_W + STEP_W * 4 + 8  # = 50+90+860+8 = 1008

    BAR_H = 44
    ROW1_Y = GRID_Y + BAR_H + 6   # AI 行
    ROW_H = 210
    ROW2_Y = ROW1_Y + ROW_H + 8   # 人手 行

    # ==========================================================
    # 3) Step バー (グレー帯 + シェブロン)
    # ==========================================================
    # 左上空白セル
    add_rect(slide, GRID_X, GRID_Y, LABEL_W, BAR_H, fill=WHITE, line=None)

    step_titles = [
        ("Step1", "商談データの確認"),
        ("Step2", "商談データの分析"),
        ("Step3", "提案戦略の検討"),
        ("Step4", "提案書の作成・提示"),
    ]
    for i, (st, ttl) in enumerate(step_titles):
        x = GRID_X + LABEL_W + STEP_W * i
        # 背景帯
        add_rect(slide, x + 2, GRID_Y, STEP_W - 4, BAR_H, fill=GRAY_BAR, line=None)
        # Step ラベル (上段、太字、ネイビー)
        add_textbox(slide, x + 2, GRID_Y + 4, STEP_W - 4, 18, st,
                    size_px=12, bold=True, color=NAVY, align='center', valign='middle')
        # タイトル (下段)
        add_textbox(slide, x + 2, GRID_Y + 22, STEP_W - 4, 20, ttl,
                    size_px=10.5, color=TEXT_BODY, align='center', valign='middle')

    # Step バー右側に「右パネル」のヘッダ帯
    add_rect(slide, PANEL_X, GRID_Y, PANEL_W, BAR_H, fill=NAVY, line=None)
    add_textbox(slide, PANEL_X, GRID_Y, PANEL_W, BAR_H,
                "Step4\n報告書の作成", size_px=11, bold=True, color=WHITE,
                align='center', valign='middle')

    # ==========================================================
    # 4) 行ラベル (AI ツール担当 / 人手担当)
    # ==========================================================
    add_rect(slide, GRID_X, ROW1_Y, LABEL_W, ROW_H, fill=GRAY_ROW, line=None)
    add_textbox(slide, GRID_X, ROW1_Y, LABEL_W, ROW_H,
                "AI ツール\n担当",
                size_px=11.5, bold=True, color=NAVY, align='center', valign='middle')

    add_rect(slide, GRID_X, ROW2_Y, LABEL_W, ROW_H, fill=GRAY_ROW, line=None)
    add_textbox(slide, GRID_X, ROW2_Y, LABEL_W, ROW_H,
                "人手\n担当",
                size_px=11.5, bold=True, color=NAVY, align='center', valign='middle')

    # ==========================================================
    # 5) AI レーン (上段)
    # ==========================================================
    # Step1 × AI: 2章 商談データ自動取込・正規化
    bx1 = GRID_X + LABEL_W + 18
    by  = ROW1_Y + 30
    bw  = STEP_W - 36
    bh  = 90
    add_box_with_text(slide, bx1, by, bw, bh,
                      "商談データの自動取込\n・正規化",
                      size_px=11, bold=True, color=NAVY)
    add_textbox(slide, bx1, by + bh, bw, 60,
                "SFA から CSV を自動取得し、\n顧客名寄せ・案件 ID 正規化を実施",
                size_px=9.5, color=GRAY_444, align='center', valign='top')
    add_num_badge(slide, bx1 + 10, by + 10, 2)

    # Step2 × AI: 3章 + 4章 (縦並び)
    bx2 = GRID_X + LABEL_W + STEP_W + 18
    add_box_with_text(slide, bx2, by - 6, bw, 70,
                      "顧客属性・\n案件履歴の分析",
                      size_px=11, bold=True, color=NAVY)
    add_num_badge(slide, bx2 + 10, by + 4, 3)

    add_box_with_text(slide, bx2, by + 76, bw, 70,
                      "類似案件の\nレコメンド生成",
                      size_px=11, bold=True, color=NAVY)
    add_num_badge(slide, bx2 + 10, by + 86, 4)
    add_textbox(slide, bx2, by + 152, bw, 40,
                "AI が過去案件 DB を参照し\n候補リストを自動生成",
                size_px=9.5, color=GRAY_444, align='center', valign='top')

    # Step3 × AI: 6章 提案ロジック候補
    bx3 = GRID_X + LABEL_W + STEP_W * 2 + 18
    add_box_with_text(slide, bx3, by, bw, bh,
                      "提案ロジック候補の\n自動生成",
                      size_px=11, bold=True, color=NAVY)
    add_textbox(slide, bx3, by + bh, bw, 60,
                "勝ちパターン DB から\nテンプレ＋論点候補を抽出",
                size_px=9.5, color=GRAY_444, align='center', valign='top')
    add_num_badge(slide, bx3 + 10, by + 10, 6)

    # Step4 × AI: 8章 提案書ドラフト
    bx4 = GRID_X + LABEL_W + STEP_W * 3 + 18
    add_box_with_text(slide, bx4, by, bw, bh,
                      "提案書ドラフトの\n自動生成",
                      size_px=11, bold=True, color=NAVY)
    add_textbox(slide, bx4, by + bh, bw, 60,
                "マクロ + テンプレで\n提案書一式（73 件 / 月）を出力",
                size_px=9.5, color=GRAY_444, align='center', valign='top')
    add_num_badge(slide, bx4 + 10, by + 10, 8)

    # AI レーン矢印: 2→3, 4→6, 6→8
    arrow_y = by + bh / 2
    add_arrow(slide, bx1 + bw, arrow_y, bx2, by + 28)              # 2→3
    add_arrow(slide, bx2 + bw, by + 110, bx3, arrow_y)              # 4→6
    add_arrow(slide, bx3 + bw, arrow_y, bx4, arrow_y)              # 6→8

    # ==========================================================
    # 6) 人手レーン (下段)
    # ==========================================================
    by2 = ROW2_Y + 30

    # Step1 × 人手: ボックスなし（営業はまだ入らない）
    add_textbox(slide, bx1, by2 + bh / 2 - 10, bw, 20,
                "（AI 側から受領）", size_px=10, color=GRAY_777, align='center', valign='middle')

    # Step2 × 人手: 5章 ヒアリング整理
    add_box_with_text(slide, bx2, by2, bw, bh,
                      "顧客 KPI ／\n提案要件のヒアリング整理",
                      size_px=11, bold=True, color=NAVY)
    add_textbox(slide, bx2, by2 + bh, bw, 60,
                "営業担当が一次商談で\n顧客の意思決定軸を整理",
                size_px=9.5, color=GRAY_444, align='center', valign='top')
    add_num_badge(slide, bx2 + 10, by2 + 10, 5)

    # Step3 × 人手: 7章 戦略選定
    add_box_with_text(slide, bx3, by2, bw, bh,
                      "戦略選定・\n差別化ポイント整理",
                      size_px=11, bold=True, color=NAVY)
    add_textbox(slide, bx3, by2 + bh, bw, 60,
                "AI 候補から営業が選定し\n顧客固有の差別化を加筆",
                size_px=9.5, color=GRAY_444, align='center', valign='top')
    add_num_badge(slide, bx3 + 10, by2 + 10, 7)

    # Step4 × 人手: 9章 最終レビュー
    add_box_with_text(slide, bx4, by2, bw, bh,
                      "最終レビュー・\nクライアント提示",
                      size_px=11, bold=True, color=NAVY)
    add_textbox(slide, bx4, by2 + bh, bw, 60,
                "提案書一式を確認のうえ\nクライアント商談で提示",
                size_px=9.5, color=GRAY_444, align='center', valign='top')
    add_num_badge(slide, bx4 + 10, by2 + 10, 9)

    # 人手レーン矢印: 5→7→9
    arrow_y2 = by2 + bh / 2
    add_arrow(slide, bx2 + bw, arrow_y2, bx3, arrow_y2)
    add_arrow(slide, bx3 + bw, arrow_y2, bx4, arrow_y2)

    # クロスレーン矢印: 8（AI）→9（人手）と 5→6（人手→AI 戦略へインプット）
    add_arrow(slide, bx4 + bw / 2, by + bh + 60, bx4 + bw / 2, by2, color=NAVY, line_w=1.4)
    add_arrow(slide, bx2 + bw / 2, by2 - 4, bx3 + bw / 4, by + bh + 4, color=GRAY_444)

    # ==========================================================
    # 7) 右パネル「本報告書の章番号と対応」
    # ==========================================================
    # AI 行に対応するセル
    panel_items_top = [
        (2, "商談データ取込の\n現状"),
        (4, "AI 分析機能の\n現状"),
    ]
    panel_items_bot = [
        (5, "ヒアリング設計の\n現状"),
        (9, "提示プロセスの\n現状"),
    ]

    # 上段パネル領域
    add_rect(slide, PANEL_X, ROW1_Y, PANEL_W, ROW_H, fill=WHITE, line=GRAY_BORDER, line_w=1.0)
    py = ROW1_Y + 12
    for n, ttl in panel_items_top:
        add_num_badge(slide, PANEL_X + 22, py + 18, n, size=24)
        add_textbox(slide, PANEL_X + 44, py, PANEL_W - 50, 40, ttl,
                    size_px=10.5, bold=True, color=NAVY, align='left', valign='middle')
        py += 56
    # 下段パネル領域
    add_rect(slide, PANEL_X, ROW2_Y, PANEL_W, ROW_H, fill=WHITE, line=GRAY_BORDER, line_w=1.0)
    py = ROW2_Y + 12
    for n, ttl in panel_items_bot:
        add_num_badge(slide, PANEL_X + 22, py + 18, n, size=24)
        add_textbox(slide, PANEL_X + 44, py, PANEL_W - 50, 40, ttl,
                    size_px=10.5, bold=True, color=NAVY, align='left', valign='middle')
        py += 56

    # メイングリッド全体を囲む細枠（左ラベル＋4Step ぶんのみ、右パネルは別枠）
    add_rect(slide, GRID_X, GRID_Y, LABEL_W + STEP_W * 4, BAR_H + ROW_H * 2 + 14,
             fill=None, line=GRAY_BORDER, line_w=1.0)

    # ==========================================================
    # 8) フッタ
    # ==========================================================
    add_textbox(slide, 44, 706, 100, 18, "6",
                size_px=10, color=GRAY_777, align='left', valign='top')
    add_textbox(slide, 100, 706, 800, 18,
                "1. 営業プロセスの実態整理 ／ 1-3. AI 提案書作成ワークフローの現状",
                size_px=9.5, color=GRAY_777, align='left', valign='top')
    add_textbox(slide, 1000, 706, 240, 18,
                "© 2026. 出典：社内営業プロセス分析資料",
                size_px=9.5, color=GRAY_777, align='right', valign='top')

    # ==========================================================
    # 保存
    # ==========================================================
    out = Path(__file__).resolve().parent.parent / "outputs" / "13-ai-sales-workflow-govthinktank.pptx"
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"saved: {out}  ({out.stat().st_size/1024:.1f} KB)")
    return out


if __name__ == "__main__":
    build()
