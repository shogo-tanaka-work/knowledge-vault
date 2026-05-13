#!/usr/bin/env python3
"""
HTML スライド一覧 → PNG → 16:9 PPTX 変換スクリプト
- 各 tmp-source/*.html を Playwright で 1280x720 PNG に変換
- python-pptx で 1スライド=1ページの PPTX を生成
"""
import sys
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

ROOT = Path(__file__).resolve().parent.parent
SLIDES_DIR = ROOT / "tmp-source"
OUTPUTS_DIR = ROOT / "outputs"
PNG_DIR = ROOT / "scripts" / ".png-cache"
PNG_DIR.mkdir(parents=True, exist_ok=True)
OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)

WIDTH, HEIGHT = 1280, 720  # 16:9 ピクセル

def render_pngs():
    files = sorted(SLIDES_DIR.glob("*.html"))
    print(f"[render] {len(files)} slides")
    out = []
    with sync_playwright() as p:
        browser = p.chromium.launch(
            executable_path="/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
        )
        ctx = browser.new_context(viewport={"width": WIDTH, "height": HEIGHT}, device_scale_factor=2)
        page = ctx.new_page()
        for f in files:
            png = PNG_DIR / f"{f.stem}.png"
            url = f"file://{f.resolve()}"
            page.goto(url, wait_until="networkidle")
            # スライド要素を厳密にクリップ
            slide_el = page.query_selector(".slide")
            if slide_el:
                box = slide_el.bounding_box()
                page.screenshot(
                    path=str(png),
                    clip={"x": box["x"], "y": box["y"], "width": box["width"], "height": box["height"]},
                    omit_background=False,
                )
            else:
                page.screenshot(path=str(png), full_page=False)
            print(f"  ✓ {png.name}")
            out.append(png)
        browser.close()
    return out

def build_pptx(pngs):
    prs = Presentation()
    # 16:9 (13.333 x 7.5 inch)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png), 0, 0,
            width=prs.slide_width, height=prs.slide_height,
        )
    out = OUTPUTS_DIR / "consulting-html-slides.pptx"
    prs.save(str(out))
    print(f"[pptx] {out} ({out.stat().st_size/1024:.1f} KB, {len(pngs)} slides)")
    return out

if __name__ == "__main__":
    pngs = render_pngs()
    if not pngs:
        print("no html files", file=sys.stderr); sys.exit(1)
    build_pptx(pngs)
